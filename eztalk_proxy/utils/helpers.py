import orjson
import re
import logging
import datetime
from typing import Any, Dict, List, Tuple, Optional
import os
import uuid

from fastapi.responses import JSONResponse
from fastapi import UploadFile, Depends
from fastapi.security import OAuth2PasswordBearer
from ..models.api_models import User
 
from ..core.config import (
    COMMON_HEADERS,
    MAX_SSE_LINE_LENGTH,
    SUPPORTED_DOCUMENT_MIME_TYPES_FOR_TEXT_EXTRACTION,
    MAX_DOCUMENT_CONTENT_CHARS_FOR_PROMPT,
)

try:
    from google.cloud import storage
    from google.auth.exceptions import DefaultCredentialsError
except ImportError:
    storage = None # type: ignore
    DefaultCredentialsError = None # type: ignore
    logging.warning(
        "google-cloud-storage or google-auth library not found. "
        "GCS upload functionality for large files (video/audio) for Gemini will not be available."
    )

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None
    logging.warning("1PyPDF2 library not found. PDF text extraction will not be available.")

try:
    import docx
except ImportError:
    docx = None
    logging.warning("python-docx library not found. DOCX text extraction will not be available.")

try:
    import olefile
except ImportError:
    olefile = None
    logging.warning("olefile library not found. DOC text extraction will not be available.")

try:
    import openpyxl
    import xlrd
except ImportError:
    openpyxl = None
    xlrd = None
    logging.warning("openpyxl/xlrd libraries not found. Excel text extraction will not be available.")

try:
    from pptx import Presentation
    pptx_available = True
except ImportError:
    pptx_available = False
    logging.warning("python-pptx library not found. PowerPoint text extraction will not be available.")

try:
    from bs4 import BeautifulSoup
    bs4_available = True
except ImportError:
    bs4_available = False
    logging.warning("BeautifulSoup4 library not found. HTML text extraction will be limited.")


logger = logging.getLogger("EzTalkProxy.Utils")


def orjson_dumps_bytes_wrapper(data: Any) -> bytes:
    return orjson.dumps(
        data,
        option=orjson.OPT_NON_STR_KEYS | orjson.OPT_PASSTHROUGH_DATETIME | orjson.OPT_APPEND_NEWLINE
    )

def error_response(
    code: int,
    msg: str,
    request_id: Optional[str] = None,
    headers: Optional[Dict[str, str]] = None
) -> JSONResponse:
    log_msg = f"错误 {code}: {msg}"
    if request_id:
        log_msg = f"RID-{request_id}: {log_msg}"
    logger.warning(log_msg)
    
    final_headers = {**COMMON_HEADERS, **(headers or {})}
    
    return JSONResponse(
        status_code=code,
        content={"error": {"message": msg, "code": code, "type": "proxy_error"}},
        headers=final_headers
    )


def extract_sse_lines(buffer: bytearray) -> Tuple[List[bytes], bytearray]:
    lines: List[bytes] = []
    start_index: int = 0
    buffer_len = len(buffer)
    while start_index < buffer_len:
        newline_index = buffer.find(b'\n', start_index)
        if newline_index == -1:
            break
        line = buffer[start_index:newline_index]
        if line.endswith(b'\r'):
            line = line[:-1]
        if len(line) > MAX_SSE_LINE_LENGTH:
            logger.warning(
                f"SSE line too long ({len(line)} bytes), exceeded MAX_SSE_LINE_LENGTH ({MAX_SSE_LINE_LENGTH}). Line skipped. "
                f"Content start: {line[:100]!r}"
            )
        else:
            lines.append(line)
        start_index = newline_index + 1
    return lines, buffer[start_index:]

def get_current_time_iso() -> str:
    return datetime.datetime.utcnow().isoformat() + "Z"

def is_gemini_2_5_model(model_name: str) -> bool:
    if not isinstance(model_name, str):
        return False
    return "gemini-2.5" in model_name.lower() or "gemini-2.5-flash-image-preview" in model_name.lower()

def _extract_text_from_pdf_pypdf2(file_path: str) -> Optional[str]:
    if not PyPDF2:
        logger.warning("Attempted to extract PDF text, but PyPDF2 library is not available.")
        return None
    text_content = ""
    try:
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            if reader.is_encrypted:
                try:
                    if reader.decrypt("") == PyPDF2.PasswordType.OWNER_PASSWORD or \
                       reader.decrypt("") == PyPDF2.PasswordType.USER_PASSWORD :
                        logger.info(f"Successfully decrypted PDF (with empty password): {file_path}")
                    else:
                        logger.warning(f"PDF file is encrypted and could not be decrypted with an empty password: {file_path}")
                        return None
                except Exception as decrypt_err:
                    logger.warning(f"Failed to decrypt PDF {file_path}: {decrypt_err}")
                    return None

            for page in reader.pages:
                try:
                    text_content += page.extract_text() or ""
                except Exception as page_extract_err:
                    logger.warning(f"Error extracting text from a page in {file_path}: {page_extract_err}")
                    continue
        return text_content.strip()
    except FileNotFoundError:
        logger.error(f"PDF file not found for extraction: {file_path}")
        return None
    except Exception as e:
        logger.error(f"Error extracting text from PDF {file_path} using PyPDF2: {e}", exc_info=True)
        return None

def _extract_text_from_docx_python_docx(file_path: str) -> Optional[str]:
    if not docx:
        logger.warning("Attempted to extract DOCX text, but python-docx library is not available.")
        return None
    try:
        doc_obj = docx.Document(file_path)
        full_text = [para.text for para in doc_obj.paragraphs]
        return "\n".join(full_text).strip()
    except FileNotFoundError:
        logger.error(f"DOCX file not found for extraction: {file_path}")
        return None
    except Exception as e:
        logger.error(f"Error extracting text from DOCX {file_path} using python-docx: {e}", exc_info=True)
        return None

def _extract_text_from_doc_olefile(file_path: str) -> Optional[str]:
    """使用olefile库从.doc文档中提取文本（简单方法）"""
    if not olefile:
        logger.warning("Attempted to extract DOC text, but olefile library is not available.")
        return None
    try:
        # 这是一个基础的.doc文本提取方法
        # 注意：.doc格式复杂，这个方法可能无法提取所有文本
        with open(file_path, 'rb') as f:
            # 尝试从.doc文件中提取可读文本
            content = f.read()
            
            # 查找可能的文本内容（简单的启发式方法）
            text_parts = []
            i = 0
            while i < len(content) - 1:
                # 查找可能的文本字符序列
                if 32 <= content[i] <= 126:  # ASCII可打印字符
                    start = i
                    while i < len(content) and 32 <= content[i] <= 126:
                        i += 1
                    if i - start > 3:  # 至少4个连续字符才认为是文本
                        text_parts.append(content[start:i].decode('ascii', errors='ignore'))
                else:
                    i += 1
            
            if text_parts:
                extracted = ' '.join(text_parts).strip()
                # 过滤掉太短或明显是垃圾的内容
                words = extracted.split()
                meaningful_words = [w for w in words if len(w) > 1 and not w.isdigit()]
                if len(meaningful_words) > 5:  # 至少5个有意义的词
                    return ' '.join(meaningful_words)
        
        return None
    except Exception as e:
        logger.error(f"Error extracting text from DOC {file_path} using olefile: {e}", exc_info=True)
        return None

def _extract_text_from_excel(file_path: str, mime_type: str) -> Optional[str]:
    """从Excel文件中提取文本"""
    try:
        text_content = []
        
        if mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            # .xlsx 文件
            if not openpyxl:
                logger.warning("openpyxl not available for .xlsx extraction")
                return None
                
            workbook = openpyxl.load_workbook(file_path)
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content.append(f"=== Sheet: {sheet_name} ===")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None:
                            row_text.append(str(cell))
                    if row_text:
                        text_content.append(" | ".join(row_text))
                        
        elif mime_type == "application/vnd.ms-excel":
            # .xls 文件
            if not xlrd:
                logger.warning("xlrd not available for .xls extraction")
                return None
                
            workbook = xlrd.open_workbook(file_path)
            for sheet_idx in range(workbook.nsheets):
                sheet = workbook.sheet_by_index(sheet_idx)
                text_content.append(f"=== Sheet: {sheet.name} ===")
                
                for row_idx in range(sheet.nrows):
                    row_text = []
                    for col_idx in range(sheet.ncols):
                        cell = sheet.cell(row_idx, col_idx)
                        if cell.value:
                            row_text.append(str(cell.value))
                    if row_text:
                        text_content.append(" | ".join(row_text))
        
        return "\n".join(text_content).strip() if text_content else None
        
    except Exception as e:
        logger.error(f"Error extracting text from Excel {file_path}: {e}", exc_info=True)
        return None

def _extract_text_from_powerpoint(file_path: str) -> Optional[str]:
    """从PowerPoint文件中提取文本"""
    if not pptx_available:
        logger.warning("python-pptx not available for PowerPoint extraction")
        return None
        
    try:
        presentation = Presentation(file_path)
        text_content = []
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            text_content.append(f"=== Slide {slide_num} ===")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())
                    
        return "\n".join(text_content).strip() if text_content else None
        
    except Exception as e:
        logger.error(f"Error extracting text from PowerPoint {file_path}: {e}", exc_info=True)
        return None

def _extract_text_from_html(file_path: str) -> Optional[str]:
    """从HTML文件中提取文本"""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            html_content = f.read()
            
        if bs4_available:
            # 使用BeautifulSoup解析HTML
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # 移除script和style标签
            for script in soup(["script", "style"]):
                script.decompose()
                
            # 提取文本
            text = soup.get_text()
            
            # 清理多余的空白
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)
            
            return text.strip() if text else None
        else:
            # 简单的HTML标签移除
            import re
            # 移除HTML标签
            text = re.sub(r'<[^>]+>', '', html_content)
            # 清理多余空白
            text = re.sub(r'\s+', ' ', text)
            return text.strip() if text else None
            
    except Exception as e:
        logger.error(f"Error extracting text from HTML {file_path}: {e}", exc_info=True)
        return None

def _extract_text_from_xml(file_path: str) -> Optional[str]:
    """从XML文件中提取文本"""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            xml_content = f.read()
            
        if bs4_available:
            soup = BeautifulSoup(xml_content, 'xml')
            text = soup.get_text()
            # 清理多余的空白
            lines = (line.strip() for line in text.splitlines())
            text = '\n'.join(line for line in lines if line)
            return text.strip() if text else None
        else:
            # 简单的XML标签移除
            import re
            text = re.sub(r'<[^>]+>', '', xml_content)
            text = re.sub(r'\s+', ' ', text)
            return text.strip() if text else None
            
    except Exception as e:
        logger.error(f"Error extracting text from XML {file_path}: {e}", exc_info=True)
        return None

def _extract_text_from_json(file_path: str) -> Optional[str]:
    """从JSON文件中提取文本内容"""
    try:
        import json
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            data = json.load(f)
            
        def extract_strings(obj, path=""):
            """递归提取JSON中的字符串值"""
            strings = []
            if isinstance(obj, dict):
                for key, value in obj.items():
                    new_path = f"{path}.{key}" if path else key
                    strings.extend(extract_strings(value, new_path))
            elif isinstance(obj, list):
                for i, item in enumerate(obj):
                    new_path = f"{path}[{i}]"
                    strings.extend(extract_strings(item, new_path))
            elif isinstance(obj, str) and obj.strip():
                strings.append(f"{path}: {obj}")
            elif obj is not None:
                strings.append(f"{path}: {str(obj)}")
            return strings
            
        text_parts = extract_strings(data)
        return "\n".join(text_parts) if text_parts else None
        
    except Exception as e:
        logger.error(f"Error extracting text from JSON {file_path}: {e}", exc_info=True)
        return None

def _extract_text_from_plain_text(file_path: str) -> Optional[str]:
    common_encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1', 'iso-8859-1']
    try:
        for encoding in common_encodings:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    return f.read().strip()
            except UnicodeDecodeError:
                logger.debug(f"Failed to decode plain text file {file_path} with encoding {encoding}")
                continue
            except FileNotFoundError:
                logger.error(f"Plain text file not found for extraction: {file_path}")
                return None
        logger.warning(f"Could not decode plain text file {file_path} with common encodings.")
        return None
    except Exception as e:
        logger.error(f"Error extracting text from plain text file {file_path}: {e}", exc_info=True)
        return None


async def extract_text_from_uploaded_document(
    uploaded_file_path: str,
    mime_type: Optional[str],
    original_filename: str
) -> Optional[str]:
    logger.info(f"Attempting to extract text from '{original_filename}' (path: {uploaded_file_path}, mime: {mime_type})")
    effective_mime_type = mime_type.lower() if mime_type else None

    if not effective_mime_type:
        logger.warning(f"No effective MIME type for '{original_filename}', cannot determine extraction method.")
        return None

    extracted_text: Optional[str] = None

    if effective_mime_type in SUPPORTED_DOCUMENT_MIME_TYPES_FOR_TEXT_EXTRACTION:
        # Microsoft Office Documents
        if effective_mime_type == "application/pdf":
            extracted_text = _extract_text_from_pdf_pypdf2(uploaded_file_path)
        elif effective_mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            extracted_text = _extract_text_from_docx_python_docx(uploaded_file_path)
        elif effective_mime_type == "application/msword":
            logger.warning(f"🔥 .doc格式文档处理：'{original_filename}' - .doc格式较老，提取效果可能不佳")
            extracted_text = _extract_text_from_doc_olefile(uploaded_file_path)
            if not extracted_text or len(extracted_text.strip()) < 10:
                extracted_text = f"""[文档解析提示]

.doc格式文档 '{original_filename}' 的内容提取遇到困难。

可能原因：
1. .doc是较老的Microsoft Word格式，结构复杂
2. 文档可能包含特殊格式或加密保护
3. 当前解析器对复杂.doc文档支持有限

建议解决方案：
1. 将文档转换为.docx格式后重新上传
2. 将文档另存为PDF格式后重新上传  
3. 复制文档内容到纯文本文件(.txt)后上传

如需帮助转换文档格式，请告知具体需求。"""
                logger.warning(f"Failed to extract meaningful content from .doc file '{original_filename}'")
        
        # Excel Documents
        elif effective_mime_type in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "application/vnd.ms-excel"]:
            extracted_text = _extract_text_from_excel(uploaded_file_path, effective_mime_type)
        
        # PowerPoint Documents  
        elif effective_mime_type in ["application/vnd.openxmlformats-officedocument.presentationml.presentation", "application/vnd.ms-powerpoint"]:
            extracted_text = _extract_text_from_powerpoint(uploaded_file_path)
        
        # Web & Markup Documents
        elif effective_mime_type == "text/html":
            extracted_text = _extract_text_from_html(uploaded_file_path)
        elif effective_mime_type in ["text/xml", "application/xml"]:
            extracted_text = _extract_text_from_xml(uploaded_file_path)
        elif effective_mime_type == "application/json":
            extracted_text = _extract_text_from_json(uploaded_file_path)
        
        # Plain Text & Other Formats
        elif effective_mime_type.startswith("text/"):
            extracted_text = _extract_text_from_plain_text(uploaded_file_path)
        else:
            logger.info(f"MIME type '{effective_mime_type}' for '{original_filename}' is in supported list but no specific extractor implemented, attempting plain text.")
            extracted_text = _extract_text_from_plain_text(uploaded_file_path)
    else:
        logger.warning(f"Unsupported MIME type for text extraction: '{effective_mime_type}' for file '{original_filename}'.")
        return None

    if extracted_text:
        if len(extracted_text) > MAX_DOCUMENT_CONTENT_CHARS_FOR_PROMPT:
            logger.info(f"Extracted text from '{original_filename}' truncated from {len(extracted_text)} to {MAX_DOCUMENT_CONTENT_CHARS_FOR_PROMPT} characters.")
            extracted_text = extracted_text[:MAX_DOCUMENT_CONTENT_CHARS_FOR_PROMPT] + \
                             f"\n[内容已截断，原始长度超过 {MAX_DOCUMENT_CONTENT_CHARS_FOR_PROMPT} 字符]"
        logger.info(f"Successfully extracted text (len: {len(extracted_text)}) from '{original_filename}'.")
        return extracted_text.strip()
    else:
        logger.warning(f"Failed to extract text from '{original_filename}' (mime: {effective_mime_type}).")
        return None

async def upload_to_gcs(
    file_obj: Any,
    original_filename: str,
    bucket_name: str,
    project_id: Optional[str] = None,
    content_type: Optional[str] = None,
    request_id: Optional[str] = None
) -> Optional[str]:
    log_prefix = f"RID-{request_id}" if request_id else "[GCS_UPLOAD]"
    
    if not storage:
        logger.error(f"{log_prefix} GCS upload skipped: google-cloud-storage library not available.")
        return None
    if not bucket_name:
        logger.error(f"{log_prefix} GCS upload skipped: GCS_BUCKET_NAME is not configured.")
        return None

    _, file_extension = os.path.splitext(original_filename)
    safe_original_filename_part = "".join(c if c.isalnum() or c in ['.', '_', '-'] else '_' for c in original_filename.rsplit('.', 1)[0])[:50]
    destination_blob_name = f"uploads/{request_id or 'unknown_req'}/{safe_original_filename_part}_{uuid.uuid4().hex[:8]}{file_extension}"

    logger.info(f"{log_prefix} Attempting to upload to GCS: bucket='{bucket_name}', blob='{destination_blob_name}'")

    try:
        if project_id:
            storage_client = storage.Client(project=project_id)
        else:
            storage_client = storage.Client()
            
        bucket = storage_client.bucket(bucket_name)
        blob = bucket.blob(destination_blob_name)

        if isinstance(file_obj, UploadFile):
            await file_obj.seek(0)
            blob.upload_from_file(file_obj.file, content_type=content_type or file_obj.content_type)
        elif hasattr(file_obj, 'read') and hasattr(file_obj, 'seek'):
            file_obj.seek(0)
            blob.upload_from_file(file_obj, content_type=content_type)
        elif isinstance(file_obj, str) and os.path.exists(file_obj):
            blob.upload_from_filename(file_obj, content_type=content_type)
        else:
            logger.error(f"{log_prefix} GCS upload failed for '{original_filename}': Invalid file_obj type or file path does not exist.")
            return None

        gcs_uri = f"gs://{bucket_name}/{destination_blob_name}"
        logger.info(f"{log_prefix} Successfully uploaded '{original_filename}' to GCS: {gcs_uri}")
        return gcs_uri
    except DefaultCredentialsError:
        logger.error(
            f"{log_prefix} GCS upload failed for '{original_filename}': Google Cloud Default Credentials not found. "
            "Ensure GOOGLE_APPLICATION_CREDENTIALS environment variable is set correctly "
            "or the runtime environment has appropriate GCS permissions.",
            exc_info=True
        )
        return None
    except Exception as e:
        logger.error(f"{log_prefix} GCS upload failed for '{original_filename}' (blob '{destination_blob_name}'): {e}", exc_info=True)
        return None

oauth2_scheme = OAuth2PasswordBearer(tokenUrl="token")

async def get_current_user(token: str = Depends(oauth2_scheme)):
   # In a real application, you would verify the token and fetch the user from a database.
   # For now, we'll just return a dummy user.
   return User(username="johndoe", email="johndoe@example.com", full_name="John Doe")